"""Generate Technical Approach & Findings slide deck for Project Caravela."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent.parent / "docs" / "executive"

# ── Colour palette ──────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
DARK_GREY = RGBColor(0x42, 0x42, 0x42)
MED_GREY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
BLUE = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
ORANGE_LIGHT = RGBColor(0xFF, 0xE0, 0xB2)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0xC8, 0xE6, 0xC9)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xE1, 0xBE, 0xE7)
TEAL = RGBColor(0x00, 0x83, 0x8F)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=13,
                     color=BLACK, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, indent) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = indent
        p.space_after = spacing
    return tf


def add_table(slide, left, top, width, height, rows, cols, data,
              header_bg=BLUE, header_fg=WHITE, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_fg
                else:
                    paragraph.font.color.rgb = BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    return table


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), BLUE)
    add_accent_bar(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), ORANGE)

    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                "Project Caravela", font_size=44, bold=True, color=BLUE)
    add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
                "Technical Approach & Findings", font_size=28, color=DARK_GREY)
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.6),
                "Brazilian E-Commerce Analytics Pipeline  |  Olist Dataset  |  ~100k orders, 2016–2018",
                font_size=16, color=MED_GREY)
    add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
                "5-Agent AI Architecture  +  Human Orchestrator",
                font_size=14, color=PURPLE, bold=True)
    add_textbox(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.5),
                "DSAI Module 2 Assignment  •  March 2026",
                font_size=12, color=MED_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Pipeline Architecture
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Pipeline Architecture", font_size=32, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                "End-to-end flow across 6 layers — modular, each independently replaceable",
                font_size=14, color=MED_GREY)

    # Pipeline flow boxes
    stages = [
        ("CSV\n9 files\n~1.1M rows", ORANGE, Inches(0.3)),
        ("Meltano\ntap-csv →\ntarget-bigquery", GREEN, Inches(2.1)),
        ("BigQuery\nolist_raw\n9 tables", BLUE, Inches(3.9)),
        ("dbt\nstaging (9 views)\nmarts (7 tables)", GREEN, Inches(5.7)),
        ("BigQuery\nolist_analytics\nStar Schema", BLUE, Inches(7.5)),
        ("Jupyter\n4 notebooks\n11 metrics", PURPLE, Inches(9.3)),
        ("Parquet\n6 files\ndata/", TEAL, Inches(11.1)),
    ]
    for label, color, left in stages:
        box = slide.shapes.add_shape(5, left, Inches(1.7), Inches(1.55), Inches(1.3))  # 5 = rounded rect
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for i, line in enumerate(label.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.font.bold = (i == 0)

    # Arrows between stages
    for i in range(len(stages) - 1):
        left = stages[i][2] + Inches(1.55)
        arr = slide.shapes.add_shape(1, left, Inches(2.15), Inches(0.55), Inches(0.06))
        arr.fill.solid()
        arr.fill.fore_color.rgb = MED_GREY
        arr.line.fill.background()

    # Streamlit box at end
    box = slide.shapes.add_shape(5, Inches(11.1), Inches(3.3), Inches(1.55), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE9, 0x1E, 0x63)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Streamlit"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "5-page dashboard"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow down from Parquet to Streamlit
    arr = slide.shapes.add_shape(1, Inches(11.8), Inches(3.0), Inches(0.06), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = MED_GREY
    arr.line.fill.background()

    # Execution boundaries
    add_accent_bar(slide, Inches(0.3), Inches(4.7), Inches(6.5), Inches(0.04), ORANGE)
    add_textbox(slide, Inches(0.3), Inches(4.8), Inches(6.5), Inches(0.4),
                "▸ AUTOMATED — Dagster-orchestrated (daily 09:00 SGT + manual trigger)",
                font_size=12, bold=True, color=ORANGE)

    add_accent_bar(slide, Inches(7.2), Inches(4.7), Inches(5.8), Inches(0.04), PURPLE)
    add_textbox(slide, Inches(7.2), Inches(4.8), Inches(5.8), Inches(0.4),
                "▸ MANUAL — Notebooks, Parquet export, Streamlit dashboard",
                font_size=12, bold=True, color=PURPLE)

    # Credentials note
    add_bullet_frame(slide, Inches(0.3), Inches(5.4), Inches(12), Inches(1.8), [
        ("All credentials from single .env file: GOOGLE_APPLICATION_CREDENTIALS, GCP_PROJECT_ID, BIGQUERY_RAW_DATASET, BIGQUERY_ANALYTICS_DATASET", False, 0),
        ("Parquet files decouple dashboard from BigQuery — Streamlit requires no GCP credentials", False, 0),
        ("dbt build interleaves models + tests — failing staging test blocks dependent marts", False, 0),
    ], font_size=12, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Tool Selection Rationale
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Tool Selection Rationale", font_size=32, bold=True, color=BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                "Each tool chosen for a specific architectural property — decoupling is the design principle",
                font_size=14, color=MED_GREY)

    tool_data = [
        ["Layer", "Selected", "Key Rationale", "Alternatives Rejected"],
        ["Ingestion", "Meltano\n(tap-csv)", "Declarative YAML; Singer protocol\nseparates extract from load;\nWRITE_TRUNCATE for full refresh", "Custom Python scripts\nbq load CLI"],
        ["Transformation", "dbt 1.11\non BigQuery", "SQL-first ELT with ref() lineage;\nstaging views + mart tables;\nall type casting centralised in staging", "pandas transforms\nBQ stored procedures"],
        ["Data Quality", "dbt-expectations\n+ singular SQL", "Integrated in DAG — failing test\nblocks downstream marts;\n76 tests, 0 failures", "Great Expectations\nSoda"],
        ["Orchestration", "Dagster 1.12\n+ dagster-dbt", "Asset-centric model; 1 asset per dbt\nmodel auto-generated;\ndaily 09:00 SGT schedule", "Airflow\nPrefect\ncron"],
        ["Analysis", "Jupyter +\ngoogle-cloud-bigquery", "Self-contained notebooks;\nParquet export decouples\ndashboard from BigQuery", "Python scripts\ndbt metrics"],
        ["Dashboard", "Streamlit 1.55", "Python-native; reads Parquet\n(no GCP credentials);\n4 global filters", "Looker, Metabase\nDash (Plotly)"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
              len(tool_data), 4, tool_data, header_bg=BLUE,
              col_widths=[Inches(1.8), Inches(2.0), Inches(4.8), Inches(3.7)])

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Star Schema Design
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ORANGE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Star Schema Design", font_size=32, bold=True, color=ORANGE)

    # Left: Facts
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.5),
                "Fact Tables", font_size=20, bold=True, color=ORANGE)
    add_bullet_frame(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), [
        ("fct_sales — order-item grain (~112k rows)", True, 0),
        ("Composite PK: order_id + order_item_id", False, 1),
        ("4 FKs → dim_customers, dim_products, dim_sellers, dim_date", False, 1),
        ("total_sale_amount = price + freight_value", False, 1),
        ("3-source CTE: items → orders → customers", False, 1),
        ("", False, 0),
        ("fct_reviews — review grain (~97k rows, post-dedup)", True, 0),
        ("PK: review_id (789 source duplicates removed)", False, 1),
        ("FK: order_id → stg_orders (NOT fct_sales)", False, 1),
        ("⚠ 756 itemless orders have reviews but no sales rows", False, 1),
        ("", False, 0),
        ("fct_payments — payment-method grain", True, 0),
        ("Compound PK: order_id + payment_sequential", False, 1),
        ("Nullable date_key via LEFT JOIN to stg_orders", False, 1),
    ], font_size=12, color=DARK_GREY)

    # Right: Dims
    add_textbox(slide, Inches(6.8), Inches(1.2), Inches(4), Inches(0.5),
                "Dimension Tables", font_size=20, bold=True, color=BLUE)
    add_bullet_frame(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(3.5), [
        ("dim_customers — PK: customer_unique_id", True, 0),
        ("Deduped via ROW_NUMBER; geolocation lat/lng (nullable)", False, 1),
        ("", False, 0),
        ("dim_products — PK: product_id", True, 0),
        ("COALESCE(english → portuguese → 'uncategorized')", False, 1),
        ("7 nullable numeric columns via SAFE_CAST", False, 1),
        ("", False, 0),
        ("dim_sellers — PK: seller_id", True, 0),
        ("Geolocation enrichment by zip code prefix", False, 1),
        ("", False, 0),
        ("dim_date — PK: date_key (DATE type)", True, 0),
        ("Generated: dbt_utils.date_spine (2016-01-01 to 2018-12-31)", False, 1),
        ("No raw data source — see ADR-001", False, 1),
    ], font_size=12, color=DARK_GREY)

    # Bottom: Why star?
    add_accent_bar(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.04), BLUE)
    add_textbox(slide, Inches(0.5), Inches(6.15), Inches(12.3), Inches(1.0),
                "Why star schema?  Single fact-to-dim joins for analytics • Dims evolve independently • "
                "Clear grain per fact table prevents 'universal fact table' anti-pattern • "
                "Staging layer (views) fixes all data defects once — every downstream consumer inherits clean data",
                font_size=13, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Data Quality Framework
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GREEN)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Data Quality Framework", font_size=32, bold=True, color=GREEN)
    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                "76 tests, 2 mechanisms, 0 failures — enforced at staging layer",
                font_size=14, color=MED_GREY)

    # Left: Generic tests
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                "Generic Tests (schema.yml)", font_size=18, bold=True, color=GREEN)
    add_bullet_frame(slide, Inches(0.5), Inches(2.0), Inches(5.5), Inches(3.0), [
        ("not_null, unique — PK/FK enforcement", False, 0),
        ("accepted_values — order_status, payment_type", False, 0),
        ("expect_column_values_to_be_between — numeric ranges", False, 0),
        ("expect_table_row_count_to_be_between — fan-out guards", False, 0),
        ("unique_combination_of_columns — compound PKs", False, 0),
        ("Temporal pair tests — timestamp ordering", False, 0),
        ("relationships — FK referential integrity", False, 0),
    ], font_size=12, color=DARK_GREY)

    # Right: Singular tests
    add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Singular Tests (SQL files)", font_size=18, bold=True, color=GREEN)
    add_bullet_frame(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.5), [
        ("assert_boleto_single_installment.sql", True, 0),
        ("Boleto payments must have exactly 1 instalment", False, 1),
        ("", False, 0),
        ("assert_payment_reconciliation.sql", True, 0),
        ("SUM(payment_value) ≈ SUM(total_sale_amount) per order", False, 1),
        ("", False, 0),
        ("assert_date_key_range.sql", True, 0),
        ("All fact table date_keys within 2016-01-01 to 2018-12-31", False, 1),
    ], font_size=12, color=DARK_GREY)

    # Bottom: Staging fixes
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
                "Staging Layer Data Defect Corrections", font_size=18, bold=True, color=GREEN)
    fix_data = [
        ["Model", "Fix", "Impact"],
        ["stg_reviews", "ROW_NUMBER dedup on review_id", "789 duplicates removed"],
        ["stg_geolocation", "Brazil bounding-box filter + AVG(lat,lng) per zip", "~1M → ~19k rows"],
        ["stg_payments", "Filter payment_type='not_defined'; clamp installments 0→1", "3 + 2 rows fixed"],
        ["stg_products", "Rename 'lenght'→'length'; COALESCE category", "610 products labelled 'uncategorized'"],
    ]
    add_table(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8),
              len(fix_data), 3, fix_data, header_bg=GREEN,
              col_widths=[Inches(2.5), Inches(5.8), Inches(4.0)])

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Market & Revenue
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ORANGE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Key Findings: Market & Revenue", font_size=32, bold=True, color=ORANGE)

    # KPI cards
    kpis = [
        ("R$15.8M", "Total GMV", ORANGE),
        ("98,353", "Orders", BLUE),
        ("R$160.51", "Avg Order Value", GREEN),
        ("97.8%", "Delivered", TEAL),
    ]
    for i, (value, label, color) in enumerate(kpis):
        left = Inches(0.5 + i * 3.1)
        box = slide.shapes.add_shape(5, left, Inches(1.2), Inches(2.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = MED_GREY
        p.alignment = PP_ALIGN.CENTER

    add_bullet_frame(slide, Inches(0.5), Inches(2.6), Inches(6), Inches(4.5), [
        ("Growth & Trajectory", True, 0),
        ("103% growth H1→H2 2017; decelerated to 38% H1 2018", False, 0),
        ("Peak: Nov 2017 — 7,451 orders (Black Friday effect)", False, 0),
        ("AOV dipped 6% during Nov 2017 (volume-driven discounting)", False, 0),
        ("", False, 0),
        ("Revenue Concentration", True, 0),
        ("Top 15 categories = 76.4% of GMV", False, 0),
        ("Category Gini 0.71, HHI 484 — diversified, no dominance", False, 0),
        ("Health & Beauty 9.1%, Watches & Gifts 8.2%, Bed/Bath/Table 7.9%", False, 0),
    ], font_size=13, color=DARK_GREY)

    add_bullet_frame(slide, Inches(6.8), Inches(2.6), Inches(6), Inches(4.5), [
        ("Payment Mix", True, 0),
        ("Credit card 77.0% — R$152 AOV", False, 0),
        ("Boleto 19.9% — R$127 AOV (20% lower than credit card)", False, 0),
        ("Credit card users: median 3 instalments, 7.3% choose 10+", False, 0),
        ("", False, 0),
        ("Freight Impact", True, 0),
        ("Freight = 14.2% of GMV (R$2.24M total)", False, 0),
        ("Christmas Supplies: 36.7% freight ratio", False, 0),
        ("Instalment promotion is a low-cost lever to boost AOV", False, 0),
    ], font_size=13, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Customers & Delivery
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Key Findings: Customers & Delivery", font_size=32, bold=True, color=BLUE)

    # KPI cards
    kpis2 = [
        ("3.1%", "Repeat Rate", RED),
        ("63.6", "NPS Proxy", BLUE),
        ("91.9%", "On-Time Delivery", GREEN),
        ("58.2%", "Hibernating", ORANGE),
    ]
    for i, (value, label, color) in enumerate(kpis2):
        left = Inches(0.5 + i * 3.1)
        box = slide.shapes.add_shape(5, left, Inches(1.2), Inches(2.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = MED_GREY
        p.alignment = PP_ALIGN.CENTER

    add_bullet_frame(slide, Inches(0.5), Inches(2.6), Inches(6), Inches(4.5), [
        ("RFM Segmentation (ref date: 2018-08-31)", True, 0),
        ("96.9% of customers made exactly 1 purchase", False, 0),
        ("3-tier frequency: F1=1, F2=2, F3=3+ orders", False, 0),
        ("Segments: Hibernating 58.2%, Promising 38.7%,", False, 0),
        ("  Loyal 1.8%, At Risk 1.0%, Champions 0.1%", False, 0),
        ("", False, 0),
        ("Actionable Insight", True, 0),
        ("'Promising' (38.7%) = recent single-purchase customers", False, 0),
        ("Highest-ROI retention target — already acquired", False, 0),
    ], font_size=13, color=DARK_GREY)

    add_bullet_frame(slide, Inches(6.8), Inches(2.6), Inches(6), Inches(4.5), [
        ("Delay-to-Satisfaction Cliff", True, 0),
        ("Early delivery:    avg score 4.30", False, 0),
        ("On-time:              avg score 3.89  (−0.41)", False, 0),
        ("1–3 days late:      avg score 3.29  (−0.60)", False, 0),
        ("4–7 days late:      avg score 2.11  (−1.19 cliff)", True, 0),
        ("7+ days late:        avg score 1.69  (−0.42)", False, 0),
        ("", False, 0),
        ("Delivery Promise", True, 0),
        ("All regions over-promise by 11–17 days", False, 0),
        ("Positive surprise drives NPS but may suppress conversion", False, 0),
    ], font_size=13, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Geographic & Sellers
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GREEN)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Key Findings: Geographic & Sellers", font_size=32, bold=True, color=GREEN)

    # Regional table
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(5), Inches(0.5),
                "Regional Performance", font_size=18, bold=True, color=GREEN)
    region_data = [
        ["Region", "GMV Share", "On-Time", "Avg Delivery", "Promise Buffer"],
        ["Southeast", "64.6%", "92.5%", "12.9 days", "11.6 days"],
        ["South", "14.5%", "92.9%", "14.0 days", "12.8 days"],
        ["Northeast", "11.9%", "85.7%", "20.6 days", "11.1 days"],
        ["Central-West", "6.5%", "92.0%", "15.5 days", "12.4 days"],
        ["North", "2.6%", "90.2%", "23.6 days", "17.5 days"],
    ]
    add_table(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(2.8),
              len(region_data), 5, region_data, header_bg=GREEN,
              col_widths=[Inches(1.6), Inches(1.3), Inches(1.3), Inches(1.6), Inches(1.7)])

    # Seller insights
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(5), Inches(0.5),
                "Seller Landscape (3,068 active)", font_size=18, bold=True, color=GREEN)
    add_bullet_frame(slide, Inches(0.5), Inches(5.3), Inches(6), Inches(2.0), [
        ("Top 18.3% of sellers = 80% of GMV", False, 0),
        ("Seller Gini 0.78 (long-tail) but HHI 35 (no monopoly)", False, 0),
        ("Platform avg seller score: 3.98/5", False, 0),
    ], font_size=13, color=DARK_GREY)

    # Seller quality tiers
    add_textbox(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.5),
                "Seller Quality Tiers (≥10 orders)", font_size=18, bold=True, color=GREEN)
    tier_data = [
        ["Tier", "Sellers", "GMV", "Avg Score", "Cancel %"],
        ["Premium", "767", "R$9.0M", "4.33", "0.0%"],
        ["Good", "364", "R$4.3M", "3.85", "1.0%"],
        ["Average", "100", "R$811K", "3.44", "2.0%"],
        ["At Risk", "37", "R$248K", "2.74", "5.4%"],
    ]
    add_table(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(2.3),
              len(tier_data), 5, tier_data, header_bg=GREEN,
              col_widths=[Inches(1.0), Inches(0.8), Inches(1.0), Inches(0.9), Inches(0.8)])

    add_bullet_frame(slide, Inches(8.3), Inches(4.3), Inches(4.5), Inches(2.5), [
        ("Key Insight", True, 0),
        ("Northeast on-time rate 85.7%", False, 0),
        ("= 7.2pp gap vs South (92.9%)", False, 0),
        ("São Paulo alone = 37.4% of total GMV", False, 0),
        ("37 At Risk sellers are a tractable", False, 0),
        ("intervention target (R$248K GMV)", False, 0),
    ], font_size=13, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Architecture Decisions & Lessons
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), DARK_GREY)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Architecture Decisions & Lessons Learned", font_size=32, bold=True, color=DARK_GREY)

    # ADR table
    adr_data = [
        ["ADR", "Decision", "Key Rationale"],
        ["ADR-001", "date_key as DATE (not INTEGER)", "dbt_utils.date_spine natively produces DATE;\nINTEGER adds unnecessary casting everywhere"],
        ["ADR-002", "Dataset names: olist_raw / olist_analytics", "'raw' is a BigQuery SQL reserved word;\nrequires backtick-quoting"],
        ["ADR-003", "fct_reviews.order_id → stg_orders\n(not fct_sales)", "756 itemless orders have reviews but\nno fct_sales rows; FK would orphan them"],
        ["ADR-004", "Meltano extractor: tap-csv", "Streaming file reads; community-maintained;\nsimpler config than tap-spreadsheets-anywhere"],
    ]
    add_table(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.8),
              len(adr_data), 3, adr_data, header_bg=DARK_GREY,
              col_widths=[Inches(1.5), Inches(4.3), Inches(6.5)])

    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
                "Non-Obvious Implementation Patterns", font_size=18, bold=True, color=DARK_GREY)
    add_bullet_frame(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.5), [
        ("3-source CTE for fct_sales", True, 0),
        ("items → orders → customers (direct join = zero matches)", False, 1),
        ("customer_id is in stg_orders; customer_unique_id in stg_customers", False, 1),
        ("", False, 0),
        ("batch_job method for Meltano loader", True, 0),
        ("gRPC idle timeout on 1M-row geolocation file", False, 1),
        ("Write API streams go idle during sequential file processing", False, 1),
    ], font_size=12, color=DARK_GREY)

    add_bullet_frame(slide, Inches(6.8), Inches(4.8), Inches(6), Inches(2.5), [
        ("@multi_asset(specs=...) for Dagster producer", True, 0),
        ("@asset(deps=...) made meltano_ingest downstream of raw tables", False, 1),
        ("@multi_asset declares it as the PRODUCER — correct topology", False, 1),
        ("", False, 0),
        ("metaplane/dbt-expectations constraint", True, 0),
        ("Fork v0.6.0 has no 'mostly' parameter on any macro", False, 1),
        ("4 proportion-based tests documented but unenforceable", False, 1),
    ], font_size=12, color=DARK_GREY)

    # ════════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — AI Multi-Agent Architecture
    # ════════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PURPLE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "AI Multi-Agent Architecture", font_size=32, bold=True, color=PURPLE)
    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                "Human Orchestrator + 5 Specialised AI Agents — contract-driven isolation",
                font_size=14, color=MED_GREY)

    # Agent table
    agent_data = [
        ["Agent", "Role", "Key Deliverables", "Worktree"],
        ["Agent 1a–1d", "Data Engineer\n(4 phases)", "Meltano config, 9 staging models,\n7 mart models, 76 dbt tests", "worktrees/agent-1"],
        ["Agent 2", "Platform Engineer", "Dagster orchestration, schedule,\nlaunch scripts, .env integration", "worktrees/agent-2"],
        ["Agent 3", "Analytics Engineer", "4 Jupyter notebooks, 6 Parquet exports,\nutils.py, analytical methodology", "worktrees/agent-3"],
        ["Agent 4", "Dashboard Engineer", "Streamlit 5-page dashboard,\n4 global filters, Lorenz/Gini charts", "worktrees/agent-4"],
        ["Agent 5", "Data Scientist", "Executive brief, slide deck,\nspeaker notes, generate_slides.py", "worktrees/agent-5"],
        ["Orchestrator", "Human + Claude\n(Opus)", "Gate management, merge conflicts,\ndiagrams, technical report, progress", "main branch"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.2),
              len(agent_data), 4, agent_data, header_bg=PURPLE,
              col_widths=[Inches(1.8), Inches(2.2), Inches(5.0), Inches(3.3)])

    # Architecture principles
    add_textbox(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
                "Architecture Principles", font_size=18, bold=True, color=PURPLE)
    add_bullet_frame(slide, Inches(0.5), Inches(5.3), Inches(5.8), Inches(2.0), [
        ("Isolation: Each agent in a separate git worktree", False, 0),
        ("Dependency gates: GATE-1→4 enforce execution order", False, 0),
        ("Contract surfaces: CLAUDE.md (spec), sources.yml", False, 0),
        ("  (Meltano↔dbt), Parquet schemas (notebooks↔dashboard)", False, 0),
    ], font_size=13, color=DARK_GREY)

    add_bullet_frame(slide, Inches(6.8), Inches(5.3), Inches(6), Inches(2.0), [
        ("42 changelog entries tracking all deviations", False, 0),
        ("Agents communicate through shared specifications,", False, 0),
        ("  not directly — the orchestrator manages gates", False, 0),
        ("  and resolves merge conflicts across worktrees", False, 0),
    ], font_size=13, color=DARK_GREY)

    # ── Save ────────────────────────────────────────────────────────────────────
    output_path = OUTPUT_DIR / "technical_approach_slides.pptx"
    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
