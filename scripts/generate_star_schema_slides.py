"""Generate Star Schema explanation slide deck for Project Caravela."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUTPUT_PATH = Path(__file__).parent.parent / "docs" / "executive" / "star_schema.pptx"

# ── Colour palette ──────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x21, 0x21, 0x21)
DARK_GREY  = RGBColor(0x42, 0x42, 0x42)
MED_GREY   = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG   = RGBColor(0xF5, 0xF5, 0xF5)
BLUE       = RGBColor(0x15, 0x65, 0xC0)
BLUE_LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
ORANGE     = RGBColor(0xE6, 0x51, 0x00)
ORANGE_LT  = RGBColor(0xFF, 0xE0, 0xB2)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LT   = RGBColor(0xC8, 0xE6, 0xC9)
RED        = RGBColor(0xC6, 0x28, 0x28)
RED_LT     = RGBColor(0xFF, 0xCD, 0xD2)
PURPLE     = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LT  = RGBColor(0xE1, 0xBE, 0xE7)
TEAL       = RGBColor(0x00, 0x83, 0x8F)
AMBER      = RGBColor(0xFF, 0x8F, 0x00)
AMBER_LT   = RGBColor(0xFF, 0xEC, 0xB3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color, line_color=None, line_pt=0, rounded=False):
    shape_id = 5 if rounded else 1
    s = slide.shapes.add_shape(shape_id, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text, size=13, bold=False,
        color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def bullets(slide, left, top, width, height, items, size=12, color=BLACK):
    """items: list of (text, bold, indent_level)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, bold, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = lvl
        p.space_after = Pt(3)
    return tf


def table(slide, left, top, width, height, data, hdr_bg=BLUE, hdr_fg=WHITE, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = "Calibri"
                para.font.bold = (r == 0)
                para.font.color.rgb = hdr_fg if r == 0 else BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_bg
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG


def header(slide, title, subtitle=None, accent=BLUE):
    box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent)
    txt(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.65),
        title, size=30, bold=True, color=accent)
    if subtitle:
        txt(slide, Inches(0.8), Inches(0.82), Inches(11.5), Inches(0.38),
            subtitle, size=13, color=MED_GREY)


def entity_box(slide, left, top, width, label, pk, fields, fill, border):
    """Draw a mini entity box for the schema diagram."""
    row_h = Inches(0.28)
    total_h = row_h * (2 + len(fields))

    # header bar
    box(slide, left, top, width, row_h, fill)
    txt(slide, left + Inches(0.08), top, width - Inches(0.1), row_h,
        label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # PK row
    pk_top = top + row_h
    box(slide, left, pk_top, width, row_h, AMBER_LT, border, 1)
    txt(slide, left + Inches(0.08), pk_top, width - Inches(0.1), row_h,
        f"🔑 {pk}", size=9, color=BLACK)

    # field rows
    for i, field in enumerate(fields):
        r_top = pk_top + row_h * (i + 1)
        bg_col = WHITE if i % 2 == 0 else LIGHT_BG
        box(slide, left, r_top, width, row_h, bg_col, border, 0.5)
        txt(slide, left + Inches(0.08), r_top, width - Inches(0.1), row_h,
            field, size=9, color=DARK_GREY)

    return total_h


def arrow_h(slide, x1, y, x2, color=MED_GREY):
    """Draw a horizontal arrow line."""
    b = slide.shapes.add_shape(1, x1, y, x2 - x1, Inches(0.04))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()


def arrow_v(slide, x, y1, y2, color=MED_GREY):
    b = slide.shapes.add_shape(1, x, y1, Inches(0.04), y2 - y1)
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: Title ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    box(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
    box(s, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), BLUE)

    # Left accent bar
    box(s, Inches(0.5), Inches(1.5), Inches(0.12), Inches(4.5), ORANGE)

    txt(s, Inches(0.8), Inches(1.5), Inches(11), Inches(1.1),
        "Star Schema Design", size=44, bold=True, color=BLUE)
    txt(s, Inches(0.8), Inches(2.65), Inches(11), Inches(0.6),
        "Project Caravela  —  olist_analytics  |  BigQuery", size=20, color=DARK_GREY)
    txt(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.5),
        "3 fact tables  ·  4 dimension tables  ·  Serving ~100k orders (2016–2018)",
        size=15, color=MED_GREY)

    # Slide index pills
    for i, (label, color) in enumerate([
        ("01  What is a Star Schema?", BLUE),
        ("02  Schema Overview", ORANGE),
        ("03  Fact Tables", GREEN),
        ("04  Dimension Tables", PURPLE),
        ("05  Design Decisions", TEAL),
    ]):
        top = Inches(4.3 + i * 0.52)
        box(s, Inches(0.8), top, Inches(0.12), Inches(0.38), color)
        txt(s, Inches(1.05), top, Inches(6), Inches(0.38), label, size=13, color=color, bold=True)

    # ── SLIDE 2: What is a Star Schema? ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    header(s, "What is a Star Schema?",
           "Facts at the centre, dimensions radiating outward — optimised for analytics joins", BLUE)

    # Left: concept diagram (simple ASCII-style boxes)
    # Central fact box
    cx = Inches(4.6)
    cy = Inches(3.4)
    fw = Inches(2.0)
    fh = Inches(1.1)
    box(s, cx, cy, fw, fh, ORANGE, ORANGE, 2, rounded=True)
    txt(s, cx, cy + Inches(0.1), fw, fh, "fct_sales\n(fact table)", size=13,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 dim boxes around it
    dims = [
        ("dim_customers", BLUE,   Inches(1.4),  Inches(2.95)),
        ("dim_products",  GREEN,  Inches(4.6),  Inches(1.15)),
        ("dim_sellers",   PURPLE, Inches(7.8),  Inches(2.95)),
        ("dim_date",      TEAL,   Inches(4.6),  Inches(5.45)),
    ]
    for label, col, dl, dt in dims:
        dw, dh = Inches(2.0), Inches(0.75)
        box(s, dl, dt, dw, dh, col, col, 2, rounded=True)
        txt(s, dl, dt + Inches(0.08), dw, dh, label, size=12,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Connector lines (approximate centres)
    mid_fact_x = cx + fw / 2
    mid_fact_y = cy + fh / 2
    # left dim → fact
    arrow_h(s, Inches(3.4), mid_fact_y, cx, BLUE)
    # top dim → fact
    arrow_v(s, mid_fact_x, Inches(1.9), cy, GREEN)
    # right dim → fact
    arrow_h(s, cx + fw, mid_fact_y, Inches(7.8), PURPLE)
    # bottom dim → fact
    arrow_v(s, mid_fact_x, cy + fh, Inches(5.45), TEAL)

    # Right: explanation bullets
    bullets(s, Inches(8.6), Inches(1.4), Inches(4.4), Inches(5.8), [
        ("Core idea", True, 0),
        ("Fact table = the event being measured", False, 1),
        ("Dimension table = descriptive context", False, 1),
        ("Each dimension is joined once per query", False, 1),
        ("", False, 0),
        ("Why it works for analytics", True, 0),
        ("Single-hop joins: fact → dim", False, 1),
        ("No multi-table chain navigation", False, 1),
        ("Dimensions evolve independently", False, 1),
        ("Clear grain prevents double-counting", False, 1),
        ("", False, 0),
        ("The 'grain' rule", True, 0),
        ("Every fact table has exactly one grain —", False, 1),
        ("the most atomic event it records", False, 1),
        ("Mixing grains in one table is the most", False, 1),
        ("common star schema mistake", False, 1),
    ], size=12, color=DARK_GREY)

    # ── SLIDE 3: Schema Overview ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    header(s, "Schema Overview",
           "olist_analytics dataset  —  3 facts × 4 dimensions  —  all relationships shown", ORANGE)

    # ── Entity boxes ──
    EW = Inches(2.35)   # entity width

    # fct_sales (centre-left)
    entity_box(s, Inches(0.25), Inches(1.3), EW, "fct_sales",
               "order_id + item_id",
               ["customer_unique_id  FK", "product_id  FK",
                "seller_id  FK", "date_key  FK",
                "price", "freight_value",
                "total_sale_amount",
                "order_delivered_…  NULLABLE",
                "order_estimated_…  NULLABLE"],
               ORANGE, ORANGE)

    # fct_reviews (centre-mid)
    entity_box(s, Inches(3.0), Inches(1.3), EW, "fct_reviews",
               "review_id",
               ["order_id  → stg_orders ⚠",
                "date_key  FK",
                "review_score",
                "review_comment_title",
                "review_comment_message"],
               ORANGE, ORANGE)

    # fct_payments (centre-right)
    entity_box(s, Inches(5.75), Inches(1.3), EW, "fct_payments",
               "order_id + payment_seq",
               ["date_key  FK  NULLABLE",
                "payment_type",
                "payment_installments",
                "payment_value"],
               ORANGE, ORANGE)

    # dim_customers (far left)
    entity_box(s, Inches(0.25), Inches(5.0), EW, "dim_customers",
               "customer_unique_id",
               ["customer_city", "customer_state",
                "customer_zip_code_prefix",
                "geolocation_lat  NULLABLE",
                "geolocation_lng  NULLABLE"],
               BLUE, BLUE)

    # dim_products (mid-left)
    entity_box(s, Inches(3.0), Inches(5.0), EW, "dim_products",
               "product_id",
               ["product_category_name_english",
                "product_name_length",
                "product_weight_g",
                "product_photos_qty"],
               GREEN, GREEN)

    # dim_sellers (mid-right)
    entity_box(s, Inches(5.75), Inches(5.0), EW, "dim_sellers",
               "seller_id",
               ["seller_city", "seller_state",
                "seller_zip_code_prefix",
                "geolocation_lat  NULLABLE",
                "geolocation_lng  NULLABLE"],
               PURPLE, PURPLE)

    # dim_date (far right)
    entity_box(s, Inches(8.5), Inches(3.4), EW, "dim_date",
               "date_key  (DATE)",
               ["year", "month", "day",
                "day_of_week", "quarter"],
               TEAL, TEAL)

    # stg_orders (referenced by fct_reviews)
    box(s, Inches(3.0), Inches(3.85), EW, Inches(0.55),
        RGBColor(0xEE, 0xEE, 0xEE), RED, 1)
    txt(s, Inches(3.05), Inches(3.87), EW - Inches(0.1), Inches(0.5),
        "stg_orders  ← fct_reviews.order_id  ⚠ ADR-003",
        size=9, color=RED)

    # FK connector lines (vertical — facts to dims below)
    # fct_sales → dim_customers
    arrow_v(s, Inches(1.1), Inches(3.5), Inches(5.0), BLUE)
    # fct_sales → dim_products
    arrow_v(s, Inches(1.4), Inches(3.5), Inches(4.85), GREEN)
    # fct_sales → dim_sellers (diagonal approximated as two lines)
    arrow_h(s, Inches(2.6), Inches(3.7), Inches(6.9), PURPLE)
    arrow_v(s, Inches(6.9), Inches(3.7), Inches(5.0), PURPLE)
    # fct_sales → dim_date
    arrow_h(s, Inches(2.6), Inches(2.1), Inches(8.6), TEAL)
    arrow_v(s, Inches(8.6), Inches(2.1), Inches(3.4), TEAL)
    # fct_reviews → dim_date
    arrow_h(s, Inches(5.35), Inches(2.3), Inches(8.7), TEAL)
    # fct_payments → dim_date
    arrow_h(s, Inches(8.1), Inches(2.5), Inches(8.75), TEAL)

    # Legend
    for i, (label, col) in enumerate([
        ("Fact table", ORANGE), ("Dimension table", BLUE),
        ("Staging ref ⚠", RED), ("FK relationship", MED_GREY),
    ]):
        lx = Inches(8.55 + (i % 2) * 2.3)
        ly = Inches(6.2 + (i // 2) * 0.45)
        box(s, lx, ly + Inches(0.08), Inches(0.22), Inches(0.22), col)
        txt(s, lx + Inches(0.3), ly, Inches(1.9), Inches(0.38), label, size=10, color=DARK_GREY)

    # ── SLIDE 4: Fact Tables ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    header(s, "Fact Tables",
           "Each fact table has one grain — the most atomic event it records", GREEN)

    tbl_data = [
        ["Fact Table", "Grain", "PK", "Rows", "Key Detail"],
        ["fct_sales", "One row per\norder line item",
         "order_id +\norder_item_id",
         "~112k",
         "3-source CTE: items→orders→customers\n"
         "total_sale_amount = price + freight_value\n"
         "Delivery timestamps repeated per item row"],
        ["fct_reviews", "One row per\ncustomer review",
         "review_id",
         "~97k\n(post-dedup)",
         "789 source duplicates removed via ROW_NUMBER\n"
         "FK → stg_orders (not fct_sales) — ADR-003\n"
         "756 itemless orders have reviews but no sales rows"],
        ["fct_payments", "One row per\npayment instalment",
         "order_id +\npayment_sequential",
         "~104k",
         "date_key from LEFT JOIN stg_orders (nullable)\n"
         "Compound PK handles split-payment orders\n"
         "payment_installments clamped 0→1 in staging"],
    ]
    table(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.0), tbl_data,
          hdr_bg=GREEN,
          col_widths=[Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.3), Inches(5.8)])

    # Why grain matters
    box(s, Inches(0.4), Inches(5.65), Inches(12.5), Inches(0.04), ORANGE)
    bullets(s, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.5), [
        ("Why grain matters:", True, 0),
        ("An order of 3 items = 3 rows in fct_sales but 1 row in fct_reviews "
         "→ always use COUNT(DISTINCT order_id) for order-level metrics", False, 1),
        ("Mixing grains (e.g. putting order-level payment totals in fct_sales) causes double-counting "
         "→ order_payment_value was deliberately excluded from fct_sales", False, 1),
    ], size=12, color=DARK_GREY)

    # ── SLIDE 5: Dimension Tables ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    header(s, "Dimension Tables",
           "Descriptive context for every fact — who, what, where, when", PURPLE)

    dim_data = [
        ["Dimension", "PK", "Key Columns", "Design Notes"],
        ["dim_customers", "customer_unique_id",
         "customer_city, state, zip\ngeolocation_lat, lng (nullable)",
         "Source has customer_id per order — deduplicated to customer_unique_id "
         "via ROW_NUMBER so one real person = one row.\n"
         "Lat/lng nullable: ~3% of customers have no matching geolocation zip."],
        ["dim_products", "product_id",
         "product_category_name_english\n7 nullable numeric attributes",
         "COALESCE(english → portuguese → 'uncategorized') resolves 610 products "
         "with no category.\nSource had misspelled columns ('lenght') — corrected "
         "in stg_products. Use product_category_name_english, never product_category_name."],
        ["dim_sellers", "seller_id",
         "seller_city, state, zip\ngeolocation_lat, lng (nullable)",
         "Geolocation enriched by zip code prefix via stg_geolocation.\n"
         "stg_geolocation bounding-box filtered to Brazil "
         "before AVG() aggregation to remove ~80 outlier coordinates."],
        ["dim_date", "date_key (DATE)",
         "year, month, day\nday_of_week, quarter",
         "Entirely generated — not loaded from any source CSV.\n"
         "dbt_utils.date_spine produces one row per day from 2016-01-01 to 2018-12-31.\n"
         "DATE type chosen over INTEGER — see ADR-001."],
    ]
    table(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5), dim_data,
          hdr_bg=PURPLE,
          col_widths=[Inches(1.9), Inches(2.0), Inches(2.8), Inches(5.8)])

    # ── SLIDE 6: Key Design Decisions ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg(s, WHITE)
    header(s, "Key Design Decisions",
           "Three non-obvious choices that shaped the schema — each with an Architecture Decision Record", TEAL)

    # Decision 1
    box(s, Inches(0.4), Inches(1.35), Inches(0.08), Inches(1.5), ORANGE)
    txt(s, Inches(0.65), Inches(1.35), Inches(12), Inches(0.45),
        "ADR-001  ·  date_key as DATE, not INTEGER", size=16, bold=True, color=ORANGE)
    bullets(s, Inches(0.65), Inches(1.8), Inches(12), Inches(0.9), [
        ("dbt_utils.date_spine natively outputs DATE — using INTEGER would require casting in every model that touches a date.", False, 0),
        ("DATE type joins natively with TIMESTAMP columns via DATE(CAST(...)) without extra functions.", False, 0),
    ], size=12, color=DARK_GREY)

    # Decision 2
    box(s, Inches(0.4), Inches(2.9), Inches(0.08), Inches(1.8), RED)
    txt(s, Inches(0.65), Inches(2.9), Inches(12), Inches(0.45),
        "ADR-003  ·  fct_reviews.order_id → stg_orders, not fct_sales", size=16, bold=True, color=RED)
    bullets(s, Inches(0.65), Inches(3.35), Inches(12), Inches(1.1), [
        ("756 orders were placed but had no line items added — they appear in stg_orders but have no rows in fct_sales.", False, 0),
        ("Pointing the FK to fct_sales would orphan those 756 review rows — referential integrity would fail.", False, 0),
        ("This is the only place in the project where a mart FK targets a staging model rather than another mart.", False, 0),
    ], size=12, color=DARK_GREY)

    # Decision 3
    box(s, Inches(0.4), Inches(4.6), Inches(0.08), Inches(2.1), GREEN)
    txt(s, Inches(0.65), Inches(4.6), Inches(12), Inches(0.45),
        "fct_sales requires a 3-source CTE  —  the customer_unique_id resolution chain", size=16, bold=True, color=GREEN)
    bullets(s, Inches(0.65), Inches(5.05), Inches(12), Inches(1.65), [
        ("stg_order_items has: order_id, product_id, seller_id  —  NO customer info", False, 0),
        ("stg_orders has: order_id → customer_id  (order-scoped, not the true PK)", False, 0),
        ("stg_customers has: customer_id → customer_unique_id  (true customer PK)", False, 0),
        ("Resolution: order_items JOIN orders (on order_id) JOIN customers (on customer_id)  "
         "→  joining items directly to customers produces zero matches", False, 0),
    ], size=12, color=DARK_GREY)

    # Code snippet box
    box(s, Inches(7.8), Inches(5.1), Inches(5.1), Inches(2.1),
        RGBColor(0x26, 0x26, 0x26))
    txt(s, Inches(7.95), Inches(5.15), Inches(4.9), Inches(2.0),
        "WITH items AS (SELECT * FROM {{ ref('stg_order_items') }}),\n"
        "     orders AS (SELECT order_id, customer_id\n"
        "                FROM {{ ref('stg_orders') }}),\n"
        "     custs  AS (SELECT customer_id, customer_unique_id\n"
        "                FROM {{ ref('stg_customers') }})\n"
        "SELECT i.*, c.customer_unique_id\n"
        "FROM items i\n"
        "JOIN orders o USING (order_id)\n"
        "JOIN custs  c USING (customer_id)",
        size=9, color=RGBColor(0xA5, 0xD6, 0xA7))

    # ── Save ─────────────────────────────────────────────────────────────
    prs.save(str(OUTPUT_PATH))
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
