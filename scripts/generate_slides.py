"""
Generate executive presentation slides (.pptx) from Parquet data.
Produces chart PNGs via plotly+kaleido, then assembles into a PowerPoint deck.
"""

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import warnings
warnings.filterwarnings("ignore")

# ── Paths ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
DATA = ROOT / "data"
ASSETS = ROOT / "docs" / "slides_assets"
ASSETS.mkdir(parents=True, exist_ok=True)
OUTPUT = ROOT / "docs" / "executive_slides.pptx"

# ── Colour palette ─────────────────────────────────────────────────────
NAVY = "#1B2A4A"
DARK_GREY = "#2C3E50"
MID_GREY = "#7F8C8D"
LIGHT_BG = "#F8F9FA"
ACCENT_BLUE = "#2980B9"
ACCENT_GREEN = "#27AE60"
ACCENT_ORANGE = "#E67E22"
ACCENT_RED = "#C0392B"
ACCENT_PURPLE = "#8E44AD"
ACCENT_TEAL = "#16A085"

REGION_COLOURS = {
    "Southeast": ACCENT_BLUE,
    "South": ACCENT_GREEN,
    "Northeast": ACCENT_ORANGE,
    "Central-West": ACCENT_PURPLE,
    "North": ACCENT_RED,
}

SEGMENT_COLOURS = {
    "Champions": "#27AE60",
    "Loyal": "#2ECC71",
    "Promising": "#3498DB",
    "At Risk": "#E67E22",
    "High Value Lost": "#C0392B",
    "Hibernating": "#95A5A6",
}

CHART_W, CHART_H = 1600, 900
CHART_SCALE = 2

# ── Load data ──────────────────────────────────────────────────────────
so = pd.read_parquet(DATA / "sales_orders.parquet")
rfm = pd.read_parquet(DATA / "customer_rfm.parquet")
sat = pd.read_parquet(DATA / "satisfaction_summary.parquet")
geo = pd.read_parquet(DATA / "geo_delivery.parquet")
sp = pd.read_parquet(DATA / "seller_performance.parquet")

# Filter to observation window
so = so[(so["year"] >= 2017) & ~((so["year"] == 2018) & (so["month"] >= 9))]
sat_filtered = sat[(sat["year"] >= 2017) & ~((sat["year"] == 2018) & (sat["month"] >= 9))]

# ── Chart generation ───────────────────────────────────────────────────
PLOT_TEMPLATE = "plotly_white"
CHART_FONT = dict(family="Arial, Helvetica, sans-serif", size=14, color=DARK_GREY)
TITLE_FONT = dict(family="Arial, Helvetica, sans-serif", size=18, color=NAVY)


def save_chart(fig, name, width=CHART_W, height=CHART_H):
    fig.update_layout(
        template=PLOT_TEMPLATE,
        font=CHART_FONT,
        title_font=TITLE_FONT,
        margin=dict(l=60, r=40, t=60, b=50),
        plot_bgcolor="white",
        paper_bgcolor="white",
    )
    path = ASSETS / f"{name}.png"
    fig.write_image(str(path), width=width, height=height, scale=CHART_SCALE)
    print(f"  Saved: {path.name}")
    return path


# ── Chart 1: Monthly GMV Trend ─────────────────────────────────────────
print("Generating charts...")

monthly = (
    so.drop_duplicates(subset=["order_id"])
    .assign(ym=lambda d: d["date_key"].dt.to_period("M").dt.to_timestamp())
    .groupby("ym")
    .agg(gmv=("total_sale_amount", "sum"), orders=("order_id", "nunique"))
    .reset_index()
)
# Use item-level GMV (sum of all items), but order count from distinct
monthly_gmv = (
    so.assign(ym=lambda d: d["date_key"].dt.to_period("M").dt.to_timestamp())
    .groupby("ym")
    .agg(gmv=("total_sale_amount", "sum"))
    .reset_index()
)
monthly_orders = (
    so.drop_duplicates(subset=["order_id"])
    .assign(ym=lambda d: d["date_key"].dt.to_period("M").dt.to_timestamp())
    .groupby("ym")["order_id"]
    .nunique()
    .reset_index(name="orders")
)
monthly = monthly_gmv.merge(monthly_orders, on="ym")

fig = make_subplots(
    rows=2, cols=1, shared_xaxes=True,
    row_heights=[0.6, 0.4], vertical_spacing=0.08,
    subplot_titles=("Monthly GMV (R$)", "Monthly Order Volume"),
)
fig.add_trace(
    go.Scatter(
        x=monthly["ym"], y=monthly["gmv"],
        fill="tozeroy", fillcolor="rgba(41,128,185,0.2)",
        line=dict(color=ACCENT_BLUE, width=2.5), name="GMV",
    ), row=1, col=1,
)
# Black Friday annotation
bf = monthly[monthly["ym"] == "2017-11-01"]
if len(bf):
    fig.add_annotation(
        x=bf["ym"].iloc[0], y=bf["gmv"].iloc[0],
        text="Black Friday<br>R$1.18M", showarrow=True,
        arrowhead=2, arrowcolor=ACCENT_RED, font=dict(size=12, color=ACCENT_RED),
        row=1, col=1,
    )
fig.add_trace(
    go.Bar(
        x=monthly["ym"], y=monthly["orders"],
        marker_color=ACCENT_BLUE, opacity=0.7, name="Orders",
    ), row=2, col=1,
)
fig.update_layout(showlegend=False, title_text="")
fig.update_yaxes(title_text="R$", row=1, col=1)
fig.update_yaxes(title_text="Orders", row=2, col=1)
save_chart(fig, "01_monthly_gmv", height=1000)

# ── Chart 2: AOV Trend ────────────────────────────────────────────────
monthly["aov"] = monthly["gmv"] / monthly["orders"]
fig = go.Figure()
fig.add_trace(go.Scatter(
    x=monthly["ym"], y=monthly["aov"],
    mode="lines+markers", line=dict(color=ACCENT_TEAL, width=2.5),
    marker=dict(size=6),
))
fig.add_hline(y=160.51, line_dash="dash", line_color=MID_GREY,
              annotation_text="Overall AOV: R$160.51", annotation_position="top left")
fig.update_layout(title="", yaxis_title="Average Order Value (R$)", xaxis_title="")
save_chart(fig, "02_aov_trend")

# ── Chart 3: Order Status Donut ────────────────────────────────────────
status = so.drop_duplicates("order_id")["order_status"].value_counts()
colours_status = {
    "delivered": ACCENT_GREEN, "shipped": ACCENT_BLUE, "canceled": ACCENT_RED,
    "invoiced": ACCENT_ORANGE, "processing": "#F39C12", "approved": MID_GREY,
    "unavailable": "#E74C3C", "created": "#BDC3C7",
}
fig = go.Figure(go.Pie(
    labels=status.index, values=status.values,
    hole=0.55, marker=dict(colors=[colours_status.get(s, MID_GREY) for s in status.index]),
    textinfo="label+percent", textposition="outside",
    textfont=dict(size=13),
))
fig.update_layout(
    title="",
    annotations=[dict(text="97.8%<br>Delivered", x=0.5, y=0.5, font_size=20,
                       font_color=ACCENT_GREEN, showarrow=False)],
    showlegend=False,
)
save_chart(fig, "03_order_status")

# ── Chart 4: RFM Segments ─────────────────────────────────────────────
seg_order = ["Champions", "Loyal", "Promising", "At Risk", "High Value Lost", "Hibernating"]
seg_counts = rfm["segment"].value_counts().reindex(seg_order)
fig = go.Figure(go.Bar(
    y=seg_counts.index, x=seg_counts.values, orientation="h",
    marker_color=[SEGMENT_COLOURS[s] for s in seg_counts.index],
    text=[f"{v:,} ({v/len(rfm)*100:.1f}%)" for v in seg_counts.values],
    textposition="outside", textfont=dict(size=13),
))
fig.update_layout(title="", xaxis_title="Number of Customers", yaxis_title="",
                  yaxis=dict(autorange="reversed"))
save_chart(fig, "04_rfm_segments")

# ── Chart 5: NPS Breakdown ────────────────────────────────────────────
nps_counts = sat_filtered["nps_category"].value_counts()
total = nps_counts.sum()
nps_colours = {"promoter": ACCENT_GREEN, "passive": "#F39C12", "detractor": ACCENT_RED}
fig = go.Figure(go.Pie(
    labels=[c.title() for c in nps_counts.index],
    values=nps_counts.values, hole=0.55,
    marker=dict(colors=[nps_colours.get(c, MID_GREY) for c in nps_counts.index]),
    textinfo="label+percent", textposition="outside", textfont=dict(size=13),
))
promoters = nps_counts.get("promoter", 0)
detractors = nps_counts.get("detractor", 0)
nps_score = (promoters - detractors) / total * 100
fig.update_layout(
    title="",
    annotations=[dict(text=f"NPS<br>{nps_score:.0f}", x=0.5, y=0.5, font_size=24,
                       font_color=ACCENT_GREEN, showarrow=False)],
    showlegend=False,
)
save_chart(fig, "05_nps_breakdown")

# ── Chart 6: Delay vs Review Score ────────────────────────────────────
bin_order = ["Early", "On-time", "1-3d late", "4-7d late", "7+d late"]
delay_data = sat_filtered[sat_filtered["delay_bin"].notna()].copy()
delay_scores = delay_data.groupby("delay_bin")["review_score"].mean().reindex(bin_order)
delay_counts = delay_data.groupby("delay_bin")["order_id"].count().reindex(bin_order)

bar_colours = [ACCENT_GREEN, ACCENT_BLUE, ACCENT_ORANGE, ACCENT_RED, "#922B21"]
fig = go.Figure(go.Bar(
    x=bin_order, y=delay_scores.values,
    marker_color=bar_colours,
    text=[f"{v:.2f}" for v in delay_scores.values],
    textposition="outside", textfont=dict(size=14, color=DARK_GREY),
))
# Annotate the cliff
fig.add_annotation(
    x="4-7d late", y=2.11, text="1.19-point cliff",
    showarrow=True, arrowhead=2, ax=-60, ay=-40,
    font=dict(size=13, color=ACCENT_RED, weight="bold"),
    arrowcolor=ACCENT_RED,
)
fig.update_layout(title="", yaxis_title="Average Review Score", xaxis_title="Delivery Timing",
                  yaxis=dict(range=[0, 5]))
save_chart(fig, "06_delay_vs_score")

# ── Chart 7: Top 10 Revenue Categories ────────────────────────────────
cat_rev = (
    so.groupby("product_category_name_english")["total_sale_amount"]
    .sum().sort_values(ascending=False).head(10)
)
fig = go.Figure(go.Bar(
    y=cat_rev.index, x=cat_rev.values, orientation="h",
    marker_color=ACCENT_BLUE,
    text=[f"R${v:,.0f}" for v in cat_rev.values],
    textposition="outside", textfont=dict(size=12),
))
fig.update_layout(title="", xaxis_title="Total Revenue (R$)", yaxis_title="",
                  yaxis=dict(autorange="reversed"))
save_chart(fig, "07_top_categories")

# ── Chart 8: Payment Method Donut ─────────────────────────────────────
pay = so.drop_duplicates("order_id")["primary_payment_type"].value_counts()
pay_colours = {
    "credit_card": ACCENT_BLUE, "boleto": ACCENT_ORANGE,
    "debit_card": ACCENT_TEAL, "voucher": ACCENT_PURPLE,
}
fig = go.Figure(go.Pie(
    labels=[l.replace("_", " ").title() for l in pay.index],
    values=pay.values, hole=0.55,
    marker=dict(colors=[pay_colours.get(p, MID_GREY) for p in pay.index]),
    textinfo="label+percent", textposition="outside", textfont=dict(size=13),
))
fig.update_layout(
    title="",
    annotations=[dict(text="77%<br>Credit Card", x=0.5, y=0.5, font_size=18,
                       font_color=ACCENT_BLUE, showarrow=False)],
    showlegend=False,
)
save_chart(fig, "08_payment_methods")

# ── Chart 9: Regional On-Time Rate ────────────────────────────────────
region_agg = geo.groupby("region").agg(
    total_orders=("total_orders", "sum"),
    on_time_orders=("on_time_orders", "sum"),
).reset_index()
region_agg["otr"] = region_agg["on_time_orders"] / region_agg["total_orders"] * 100
region_agg = region_agg.sort_values("otr", ascending=True)

fig = go.Figure(go.Bar(
    y=region_agg["region"], x=region_agg["otr"], orientation="h",
    marker_color=[REGION_COLOURS.get(r, MID_GREY) for r in region_agg["region"]],
    text=[f"{v:.1f}%" for v in region_agg["otr"]],
    textposition="outside", textfont=dict(size=14),
))
fig.add_vline(x=91.9, line_dash="dash", line_color=MID_GREY,
              annotation_text="National: 91.9%", annotation_position="top")
fig.update_layout(title="", xaxis_title="On-Time Delivery Rate (%)", yaxis_title="",
                  xaxis=dict(range=[80, 96]))
save_chart(fig, "09_regional_ontime")

# ── Chart 10: Regional GMV ────────────────────────────────────────────
rgmv = (
    so.drop_duplicates("order_id")
    .groupby("customer_region")["total_sale_amount"]
    .sum().sort_values(ascending=True).reset_index()
)
rgmv.columns = ["region", "gmv"]
fig = go.Figure(go.Bar(
    y=rgmv["region"], x=rgmv["gmv"], orientation="h",
    marker_color=[REGION_COLOURS.get(r, MID_GREY) for r in rgmv["region"]],
    text=[f"R${v/1e6:.1f}M ({v/rgmv['gmv'].sum()*100:.1f}%)" for v in rgmv["gmv"]],
    textposition="outside", textfont=dict(size=13),
))
fig.update_layout(title="", xaxis_title="Gross Merchandise Value (R$)", yaxis_title="")
save_chart(fig, "10_regional_gmv")

# ── Chart 11: Seller Scatter ──────────────────────────────────────────
sp_active = sp[(sp["order_count"] >= 10) & sp["avg_review_score"].notna()].copy()
fig = px.scatter(
    sp_active, x="gmv", y="avg_review_score", size="order_count",
    color="seller_region", color_discrete_map=REGION_COLOURS,
    opacity=0.6, size_max=25,
    labels={"gmv": "GMV (R$)", "avg_review_score": "Avg Review Score",
            "order_count": "Orders", "seller_region": "Region"},
)
fig.add_hline(y=3.5, line_dash="dash", line_color=ACCENT_RED,
              annotation_text="Quality threshold", annotation_position="top left")
fig.update_layout(title="", legend_title="Region")
save_chart(fig, "11_seller_scatter")

print("\nAll charts generated. Building PPTX...")

# ══════════════════════════════════════════════════════════════════════
# PPTX GENERATION
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY_RGB = RGBColor(0x1B, 0x2A, 0x4A)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE_RGB = RGBColor(0x29, 0x80, 0xB9)
ACCENT_GREEN_RGB = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE_RGB = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_RED_RGB = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_GREY_RGB = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GREY_RGB = RGBColor(0x2C, 0x3E, 0x50)
MID_GREY_RGB = RGBColor(0x7F, 0x8C, 0x8D)


def add_bg(slide, colour=NAVY_RGB):
    """Set slide background colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_shape_bg(slide, left, top, width, height, colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_colour=DARK_GREY_RGB, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_colour
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_kpi_card(slide, left, top, value, label, colour=ACCENT_BLUE_RGB):
    """Add a KPI card with big number and label."""
    card_w, card_h = Inches(2.8), Inches(1.6)
    shape = add_shape_bg(slide, left, top, card_w, card_h, WHITE_RGB)
    shape.shadow.inherit = False

    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15),
                 card_w - Inches(0.3), Inches(0.8),
                 value, font_size=28, font_colour=colour, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.9),
                 card_w - Inches(0.3), Inches(0.5),
                 label, font_size=13, font_colour=MID_GREY_RGB,
                 alignment=PP_ALIGN.CENTER)


def add_chart_slide(slide, title, subtitle, chart_path, bullets=None):
    """Standard chart slide: header bar + chart + optional bullets."""
    # Header bar
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY_RGB)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=28, font_colour=WHITE_RGB, bold=True)
    add_text_box(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
                 subtitle, font_size=14, font_colour=RGBColor(0xBD, 0xC3, 0xC7))

    # Chart image
    if bullets:
        img_left, img_w = Inches(0.3), Inches(8.5)
        slide.shapes.add_picture(
            str(chart_path), img_left, Inches(1.3), img_w, Inches(5.5))
        # Bullets panel
        txBox = slide.shapes.add_textbox(
            Inches(9.0), Inches(1.5), Inches(4.0), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GREY_RGB
            p.font.name = "Arial"
            p.space_after = Pt(8)
            p.level = 0
    else:
        img_left, img_w = Inches(0.8), Inches(11.5)
        slide.shapes.add_picture(
            str(chart_path), img_left, Inches(1.3), img_w, Inches(5.8))


def add_text_slide(slide, title, content_blocks):
    """Text-heavy slide with header bar + content blocks."""
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY_RGB)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=28, font_colour=WHITE_RGB, bold=True)

    y = Inches(1.4)
    for block_title, block_bullets in content_blocks:
        add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.4),
                     block_title, font_size=18, font_colour=ACCENT_BLUE_RGB, bold=True)
        y += Inches(0.45)
        for bullet in block_bullets:
            add_text_box(slide, Inches(1.1), y, Inches(11), Inches(0.35),
                         f"\u2022  {bullet}", font_size=13, font_colour=DARK_GREY_RGB)
            y += Inches(0.35)
        y += Inches(0.15)


# ── Slide 1: Title ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY_RGB)

# Accent bar
add_shape_bg(slide, Inches(0.8), Inches(2.0), Inches(1.0), Inches(0.06), ACCENT_BLUE_RGB)

add_text_box(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(1.0),
             "Brazilian E-Commerce Insights", font_size=42, font_colour=WHITE_RGB, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(10), Inches(0.6),
             "Project Caravela \u2014 Data Pipeline Findings", font_size=22,
             font_colour=RGBColor(0xBD, 0xC3, 0xC7))
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.5),
             "Analysis Period: January 2017 \u2013 August 2018  |  ~100,000 Orders  |  R$15.8M GMV",
             font_size=16, font_colour=RGBColor(0x95, 0xA5, 0xA6))

# Bottom accent
add_shape_bg(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), ACCENT_BLUE_RGB)

# ── Slide 2: Executive Summary (KPIs) ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY_RGB)
add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
             "Executive Summary", font_size=32, font_colour=WHITE_RGB, bold=True)

# KPI cards row
add_kpi_card(slide, Inches(0.5), Inches(1.4), "R$15.8M", "Total GMV", ACCENT_BLUE_RGB)
add_kpi_card(slide, Inches(3.6), Inches(1.4), "98,353", "Total Orders", ACCENT_BLUE_RGB)
add_kpi_card(slide, Inches(6.7), Inches(1.4), "63.6", "NPS Score", ACCENT_GREEN_RGB)
add_kpi_card(slide, Inches(9.8), Inches(1.4), "91.9%", "On-Time Rate", ACCENT_GREEN_RGB)

# Second row
add_kpi_card(slide, Inches(0.5), Inches(3.3), "R$160.51", "Avg Order Value", ACCENT_BLUE_RGB)
add_kpi_card(slide, Inches(3.6), Inches(3.3), "3.1%", "Repeat Purchase Rate", ACCENT_ORANGE_RGB)
add_kpi_card(slide, Inches(6.7), Inches(3.3), "0.46%", "Cancellation Rate", ACCENT_GREEN_RGB)
add_kpi_card(slide, Inches(9.8), Inches(3.3), "~103%", "H1\u2192H2 2017 Growth", ACCENT_BLUE_RGB)

# Key bullets
bullets_text = [
    "\u2022  GMV doubled from H1 to H2 2017, driven by volume (not higher order values)",
    "\u2022  96.9% of customers make only one purchase \u2014 retention is the critical challenge",
    "\u2022  Northeast delivery lags at 85.7% on-time (vs 92.9% South) \u2014 impacting satisfaction",
    "\u2022  Southeast dominates with 64.6% of GMV; S\u00e3o Paulo alone = 37.4%",
]
y = Inches(5.2)
for bt in bullets_text:
    add_text_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.35),
                 bt, font_size=14, font_colour=DARK_GREY_RGB)
    y += Inches(0.38)

# ── Slide 3: Monthly GMV Trend ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Market Overview: Monthly GMV & Volume",
                "Jan 2017 \u2013 Aug 2018  |  Peak: Nov 2017 (Black Friday)",
                ASSETS / "01_monthly_gmv.png",
                bullets=[
                    "H1\u2192H2 2017: ~103% GMV growth",
                    "H2 2017\u2192H1 2018: ~38% (decelerating)",
                    "Nov 2017 peak: 7,451 orders, R$1.18M",
                    "Black Friday was volume-driven \u2014 AOV dropped 6%",
                    "Sep/Oct 2018 excluded (data artefacts)",
                ])

# ── Slide 4: AOV Trend ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Average Order Value Trend",
                "Stable at R$160.51 \u2014 growth is volume-driven, not ticket-size-driven",
                ASSETS / "02_aov_trend.png",
                bullets=[
                    "Overall AOV: R$160.51",
                    "Remarkably flat across 20 months",
                    "Credit card AOV: R$152",
                    "Boleto AOV: R$127 (\u221220%)",
                    "Instalment access enables larger baskets",
                ])

# ── Slide 5: Order Fulfilment ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Order Fulfilment & Cancellation",
                "97.8% delivered  |  0.46% cancellation rate",
                ASSETS / "03_order_status.png",
                bullets=[
                    "97.8% of orders delivered successfully",
                    "Cancellation rate: just 0.46% (448 orders)",
                    "Strong inventory management",
                    "Remaining 1.7%: in-transit pipeline statuses",
                ])

# ── Slide 6: RFM Segments ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Customer Segmentation (RFM Analysis)",
                "95,420 unique customers  |  Reference date: 31 Aug 2018",
                ASSETS / "04_rfm_segments.png",
                bullets=[
                    "58.2% Hibernating (55,552)",
                    "38.7% Promising (36,956)",
                    "Only 3.1% repeat purchase rate",
                    "131 Champions (0.1%) \u2014 crown jewels",
                    "Marketplace model = low platform loyalty",
                    "Promising segment = top retention target",
                ])

# ── Slide 7: NPS ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Customer Satisfaction \u2014 NPS Proxy",
                "Net Promoter Score: 63.6 (Strong)",
                ASSETS / "05_nps_breakdown.png",
                bullets=[
                    "Promoters (score 4\u20135): 77.7%",
                    "Detractors (score 1\u20132): 14.1%",
                    "NPS = 77.7% \u2212 14.1% = 63.6",
                    "Strong overall, but delivery delays",
                    "are the #1 driver of detractors",
                ])

# ── Slide 8: Delay vs Score (KEY INSIGHT) ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Key Insight: Delivery Delays Destroy Satisfaction",
                "2.61-point score drop from early delivery to 7+ days late",
                ASSETS / "06_delay_vs_score.png",
                bullets=[
                    "Early delivery: score 4.30",
                    "On-time: 3.89",
                    "1\u20133 days late: 3.29",
                    "4\u20137 days late: 2.11 (\u22121.19 cliff!)",
                    "7+ days late: 1.69",
                    "",
                    "The 4-day mark is where",
                    "customer patience breaks.",
                    "6,359 orders (6.6%) were late.",
                ])

# ── Slide 9: Top Categories ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Product & Revenue: Top Categories",
                "Top 15 categories = 76.4% of GMV  |  74 categories total",
                ASSETS / "07_top_categories.png",
                bullets=[
                    "#1 Health & Beauty: 9.1%",
                    "#2 Watches & Gifts: 8.2%",
                    "#3 Bed, Bath & Table: 7.9%",
                    "",
                    "Category Gini: 0.71 (long tail)",
                    "HHI: 484 (competitive, no dominance)",
                    "Freight = 14.2% of GMV (R$2.24M)",
                ])

# ── Slide 10: Payment Methods ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Payment Method Distribution",
                "Credit card dominance with instalment culture",
                ASSETS / "08_payment_methods.png",
                bullets=[
                    "Credit Card: 77.0% of orders",
                    "Boleto: 19.9%",
                    "Debit: 1.5% | Voucher: 1.6%",
                    "",
                    "Median instalments: 3",
                    "7.3% choose 10+ instalments",
                    "CC AOV 20% higher than boleto",
                ])

# ── Slide 11: Regional Delivery ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Delivery Performance by Region",
                "National on-time rate: 91.9%  |  Northeast lags at 85.7%",
                ASSETS / "09_regional_ontime.png",
                bullets=[
                    "South: 92.9% (best)",
                    "Southeast: 92.5%",
                    "Central-West: 92.0%",
                    "North: 90.2%",
                    "Northeast: 85.7% (worst)",
                    "",
                    "NE transit: 20.6 days avg",
                    "vs SE: 12.9 days avg",
                    "7.2pp gap \u2192 satisfaction gap",
                ])

# ── Slide 12: Regional GMV ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Geographic Revenue Distribution",
                "Southeast = 64.6% of GMV  |  S\u00e3o Paulo alone = 37.4%",
                ASSETS / "10_regional_gmv.png",
                bullets=[
                    "Southeast: R$10.2M (64.6%)",
                    "South: R$2.3M (14.5%)",
                    "Northeast: R$1.9M (11.9%)",
                    "Central-West: R$1.0M (6.5%)",
                    "North: R$0.4M (2.6%)",
                    "",
                    "North + Central-West = 9.1%",
                    "but >20% of Brazil\u2019s population",
                ])

# ── Slide 13: Seller Landscape ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_chart_slide(slide, "Seller Performance Landscape",
                "3,068 sellers  |  Top 18.3% generate 80% of GMV",
                ASSETS / "11_seller_scatter.png",
                bullets=[
                    "Seller Gini: 0.78 (high inequality)",
                    "CR4: 6.1% (no monopoly)",
                    "HHI: 35 (competitive market)",
                    "",
                    "767 Premium sellers (score \u22654.3)",
                    "37 At Risk sellers (score <3.0)",
                    "Avg platform score: 3.98/5",
                ])

# ── Slide 14: Key Risks ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_text_slide(slide, "Key Risks & Mitigations", [
    ("Technical Risk: Historical Data Limitations", [
        "Dataset covers 2016\u20132018 \u2014 consumer behaviour has evolved since then",
        "Mitigation: Pipeline is repeatable; refresh with current data and monitor metric drift quarterly",
    ]),
    ("Business Risk: Single-Purchase Customer Dominance", [
        "96.9% of customers buy only once \u2014 CAC amortised over a single transaction",
        "As high-penetration regions saturate, acquisition costs increase unsustainably",
        "Mitigation: Target 36,956 Promising customers with post-purchase engagement programmes",
    ]),
    ("Business Risk: Regional Delivery Disparity", [
        "Northeast on-time rate 85.7% vs South 92.9% \u2014 a 7.2pp gap",
        "Late deliveries score 1.69 vs 4.30 for early \u2014 compounding NE growth barriers",
        "Mitigation: Regional carrier partnerships + distribution centres in NE and North",
    ]),
])

# ── Slide 15: Strategic Recommendations ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE_RGB)
add_text_slide(slide, "Strategic Recommendations", [
    ("1. Customer Retention Programme", [
        "Target 36,956 Promising segment customers within 30 days of first purchase",
        "Evidence: Repeat rate declining from 5.3% (2017) to 2.6% (2018 cohorts)",
    ]),
    ("2. North & Northeast Delivery Infrastructure", [
        "Partner with regional carriers; establish fulfilment hubs in Recife and Manaus",
        "Evidence: NE on-time rate 85.7% vs South 92.9%; North transit 23.6 days avg",
    ]),
    ("3. Optimise Delivery Promise Accuracy", [
        "A/B test tighter estimates in SE (actual 12.9d vs promised 24.5d = 11.6d buffer)",
        "Hypothesis: Tighter promises increase checkout conversion without hurting on-time rate",
    ]),
    ("4. Seller Quality Intervention", [
        "37 At Risk sellers (avg score 2.74, 5.4% cancel rate) \u2014 90-day improvement window",
    ]),
    ("5. Instalment-Driven AOV Uplift", [
        "CC AOV R$152 vs boleto R$127; promote instalments to boleto-inclined buyers",
    ]),
])

# ── Slide 16: Thank You ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_RGB)
add_shape_bg(slide, Inches(0.8), Inches(2.8), Inches(1.0), Inches(0.06), ACCENT_BLUE_RGB)
add_text_box(slide, Inches(0.8), Inches(3.1), Inches(10), Inches(0.8),
             "Thank You", font_size=42, font_colour=WHITE_RGB, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(0.6),
             "Project Caravela  |  Questions & Discussion", font_size=20,
             font_colour=RGBColor(0xBD, 0xC3, 0xC7))
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(0.5),
             "Interactive dashboard available at: streamlit run dashboard.py",
             font_size=14, font_colour=RGBColor(0x95, 0xA5, 0xA6))
add_shape_bg(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), ACCENT_BLUE_RGB)

# ── Save ──────────────────────────────────────────────────────────────
prs.save(str(OUTPUT))
print(f"\nPresentation saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
